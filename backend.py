from fastapi import FastAPI, UploadFile, File
from fastapi.responses import FileResponse
from pydantic import BaseModel
from fastapi.middleware.cors import CORSMiddleware
from sentence_transformers import SentenceTransformer
import google.generativeai as genai
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl
import faiss
import numpy as np
from io import BytesIO
from dotenv import load_dotenv
import os

# ------------------ Load Environment ------------------
load_dotenv()
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))

# ------------------ FastAPI Init ------------------
app = FastAPI()
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ------------------ Embedding & LLM ------------------
embedding_model = SentenceTransformer("all-MiniLM-L6-v2")
llm_model = genai.GenerativeModel("gemini-1.5-flash")

# ------------------ Global Storage ------------------
file_chunks = []
index = None
uploaded_ext = None

# ------------------ Utility Functions ------------------
def chunk_text(text, chunk_size=500):
    words = text.split()
    return [" ".join(words[i:i + chunk_size]) for i in range(0, len(words), chunk_size)]

def extract_text_from_pdf(file):
    pdf_reader = PdfReader(file)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text() or ""
    return text

def extract_text_from_docx(file):
    # Wrap SpooledTemporaryFile into BytesIO
    file_bytes = BytesIO(file.read())
    doc = DocxDocument(file_bytes)
    text = " ".join([para.text for para in doc.paragraphs])
    file.seek(0)  # reset pointer for reuse
    return text

def extract_text_from_excel(file):
    # Convert to BytesIO to avoid SpooledTemporaryFile issue
    file_bytes = BytesIO(file.read())
    workbook = openpyxl.load_workbook(file_bytes, data_only=True)
    text = ""
    for sheet in workbook.sheetnames:
        ws = workbook[sheet]
        for row in ws.iter_rows(values_only=True):
            row_text = " ".join([str(cell) for cell in row if cell])
            text += row_text + " "
    file.seek(0)
    return text

def extract_text_from_ppt(file):
    # Convert to BytesIO to avoid SpooledTemporaryFile issue
    file_bytes = BytesIO(file.read())
    prs = Presentation(file_bytes)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
    file.seek(0)
    return text

# ------------------ API Models ------------------
class QARequest(BaseModel):
    query: str

# ------------------ Routes ------------------
@app.post("/upload-file/")
async def upload_file(file: UploadFile = File(...)):
    global file_chunks, index, uploaded_ext

    ext = file.filename.split(".")[-1].lower()
    uploaded_ext = ext

    # Extract text based on file type
    if ext == "pdf":
        text = extract_text_from_pdf(file.file)
    elif ext in ["docx", "doc"]:
        text = extract_text_from_docx(file.file)
    elif ext in ["xlsx", "xls"]:
        text = extract_text_from_excel(file.file)
    elif ext in ["pptx", "ppt"]:
        text = extract_text_from_ppt(file.file)
    else:
        return {"error": "Unsupported file format"}

    # Chunk text & create embeddings
    file_chunks = chunk_text(text)
    embeddings = embedding_model.encode(file_chunks)
    dim = embeddings.shape[1]
    index = faiss.IndexFlatL2(dim)
    index.add(np.array(embeddings))

    return {"message": f"{ext.upper()} uploaded and processed successfully", "chunks": len(file_chunks)}

@app.post("/ask")
async def ask_question(request: QARequest):
    if not file_chunks:
        return {"answer": "⚠️ Please upload a file first."}

    query = request.query.lower()

    # Special case: word count
    if "word count" in query or "total words" in query or "how many words" in query:
        all_text = " ".join(file_chunks)
        word_count = len(all_text.split())
        return {"answer": f"The document has approximately {word_count} words."}

    # Normal QnA
    query_embedding = embedding_model.encode([request.query])
    D, I = index.search(query_embedding, k=3)
    context = " ".join([file_chunks[i] for i in I[0]])

    prompt = f"""
    Context: {context}
    Question: {request.query}
    Answer only from context. 
    If not found, say: "I couldn't find that information in the document."
    """
    response = llm_model.generate_content(prompt)
    return {"answer": response.text}

@app.get("/convert-to-word/")
async def convert_to_word():
    if not file_chunks:
        return {"error": "⚠️ Please upload a file first."}

    doc = DocxDocument()
    doc.add_heading("Converted Document Content", 0)
    for chunk in file_chunks:
        doc.add_paragraph(chunk)

    file_path = "converted.docx"
    doc.save(file_path)

    return FileResponse(
        file_path,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        filename="converted.docx"
    )
